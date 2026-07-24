from __future__ import annotations

import asyncio
from pathlib import Path


class PptxParser:
    name = "pptx"

    async def parse(
        self,
        filename: str | Path,
        *,
        as_json: bool = False,
    ) -> str | dict[str, object]:
        slides = await asyncio.to_thread(_read_presentation, Path(filename))
        if as_json:
            return {"slides": slides}
        return "\n\n".join(_slide_markdown(slide) for slide in slides)


def _read_presentation(path: Path) -> list[dict[str, object]]:
    from pptx import Presentation

    presentation = Presentation(path)
    slides: list[dict[str, object]] = []
    has_content = False
    for position, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        elements: list[dict[str, object]] = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                elements.append({"type": "table", "rows": rows})
                has_content = has_content or any(any(row) for row in rows)
            elif getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        elements.append(
                            {
                                "type": "paragraph",
                                "level": paragraph.level,
                                "text": text,
                            }
                        )
                        has_content = True
        has_content = has_content or bool(title)
        slides.append(
            {
                "position": position,
                "title": title,
                "elements": elements,
            }
        )
    if slides and not has_content:
        raise ValueError("PowerPoint contains no extractable content.")
    return slides


def _slide_markdown(slide: dict[str, object]) -> str:
    title = str(slide["title"])
    heading = f"# 幻灯片 {slide['position']}" + (f"：{title}" if title else "")
    lines = [heading]
    for element in slide["elements"]:  # type: ignore[union-attr]
        if element["type"] == "table":
            rows = element["rows"]
            if rows:
                lines.extend(["", *_table_markdown(rows)])
        else:
            level = int(element.get("level", 0))
            lines.append(f"{'  ' * level}{element['text']}")
    return "\n".join(lines)


def _table_markdown(rows: list[list[str]]) -> list[str]:
    width = max(len(row) for row in rows)
    padded = [row + [""] * (width - len(row)) for row in rows]

    def line(values: list[str]) -> str:
        return "| " + " | ".join(value.replace("|", "\\|") for value in values) + " |"

    return [line(padded[0]), line(["---"] * width), *(line(row) for row in padded[1:])]
